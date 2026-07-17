import os
import sys

# Automatically install python-pptx if it is not installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx library...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Colors
    bg_color = RGBColor(3, 7, 18)        # Dark Indigo-Gray
    text_color = RGBColor(245, 245, 245) # Light Off-White
    muted_color = RGBColor(161, 161, 170)# Gray
    indigo_color = RGBColor(99, 102, 241)# Indigo Accent
    purple_color = RGBColor(139, 92, 246)# Purple Accent
    emerald_color = RGBColor(16, 185, 129)# Emerald Green
    red_color = RGBColor(239, 68, 68)    # Light Red
    amber_color = RGBColor(245, 158, 11)  # Amber Yellow

    def apply_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_tag(slide, text):
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.5))
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"▪  {text.upper()}"
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = indigo_color

    # ═══════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ═══════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide1)
    
    add_tag(slide1, "Runtime Security for Autonomous AI")
    
    # Title text box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The security layer\nyour AI agents are missing."
    p.font.name = 'Arial'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Subtitle text box
    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(10), Inches(1.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Kavach AI is a zero-trust runtime firewall that intercepts, analyzes, and secures every autonomous agent action — before it reaches your infrastructure."
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = muted_color
    
    # Info footer
    info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(5), Inches(0.5))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "🛡️ kavach-ai.dev  |  Hackathon 2026"
    p_info.font.name = 'Arial'
    p_info.font.size = Pt(12)
    p_info.font.color.rgb = purple_color

    # ═══════════════════════════════════════════════
    # SLIDE 2: THE PROBLEM
    # ═══════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide2)
    
    add_tag(slide2, "The Problem")
    
    title_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI agents are powerful, but dangerously unguarded."
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # 4 threat boxes as shapes
    threats = [
        ("Prompt Injection & Jailbreaks", "Dynamic data sources poison agent context to hijack model behavior and override safety rules.", red_color),
        ("Data Exfiltration", "Agents read sensitive databases and silently upload credentials or PII to external endpoints.", amber_color),
        ("Privilege Escalation", "Agents grant themselves admin roles, create new accounts, or modify IAM configurations.", amber_color),
        ("Destructive Operations", "Unguarded tool calls execute destructive commands, drop tables, or wipe entire filesystems.", red_color)
    ]
    
    positions = [
        (Inches(0.8), Inches(2.5)), # Top Left
        (Inches(6.8), Inches(2.5)), # Top Right
        (Inches(0.8), Inches(4.7)), # Bottom Left
        (Inches(6.8), Inches(4.7))  # Bottom Right
    ]
    
    for i, (title, desc, color) in enumerate(threats):
        x, y = positions[i]
        
        # Draw background shape
        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        shape.line.color.rgb = RGBColor(40, 40, 50)
        shape.line.width = Pt(1)
        
        # Add Text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = muted_color
        p2.space_before = Pt(8)

    # ═══════════════════════════════════════════════
    # SLIDE 3: THE SOLUTION
    # ═══════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide3)
    
    add_tag(slide3, "Our Solution")
    
    title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Kavach AI: Zero-Trust Interception Proxy"
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Diagram box (Agent -> Proxy -> Server)
    nodes = ["🤖 AI Agent", "🛡️ Kavach Proxy", "🗄️ MCP Servers"]
    colors = [muted_color, purple_color, muted_color]
    
    for i, name in enumerate(nodes):
        x = Inches(0.8 + i * 4.0)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(3.2), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        if name == "🛡️ Kavach Proxy":
            shape.line.color.rgb = purple_color
            shape.line.width = Pt(2)
        else:
            shape.line.color.rgb = RGBColor(40, 40, 50)
            shape.line.width = Pt(1)
            
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = text_color
        
        # Draw Arrow
        if i < 2:
            arrow_box = slide3.shapes.add_textbox(x + Inches(3.2), Inches(2.55), Inches(0.8), Inches(0.5))
            tf_a = arrow_box.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.alignment = PP_ALIGN.CENTER
            p_a.text = "➔"
            p_a.font.size = Pt(24)
            p_a.font.color.rgb = purple_color
            
    # Solution descriptions
    sol_boxes = [
        ("🔍 Intercept & Analyze", "Intent classification, injection scanning, and semantic threat analysis on every call in real-time."),
        ("📊 Score & Decide", "Multi-factor risk scoring, behavior profiling, and configurable policy engine rule checks."),
        ("🧑‍💻 Human Approval", "High-risk escalation operations halt execution, waiting for security analyst approval.")
    ]
    
    for i, (title, desc) in enumerate(sol_boxes):
        x = Inches(0.8 + i * 4.0)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.8), Inches(3.7), Inches(2.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        shape.line.color.rgb = RGBColor(40, 40, 50)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = indigo_color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = muted_color
        p2.space_before = Pt(8)

    # ═══════════════════════════════════════════════
    # SLIDE 4: HOW IT WORKS (9 STAGES)
    # ═══════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide4)
    
    add_tag(slide4, "How it Works")
    
    title_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "9-Stage Security Pipeline"
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    stages = [
        ("01", "MCP Intercept", "Raw tool call interception"),
        ("02", "Intent Analysis", "Semantic threat check"),
        ("03", "Injection Detect", "Pattern override scans"),
        ("04", "Behavior Profile", "Flag anomalous arguments"),
        ("05", "Risk Scoring", "Dynamic factor score"),
        ("06", "Trust Engine", "Violations trust decay"),
        ("07", "Policy Engine", "Configurable rule tests"),
        ("08", "Human Approval", "HALT on critical alerts"),
        ("09", "Execution Decis.", "Block or forward decision")
    ]
    
    for i, (num, name, desc) in enumerate(stages):
        row = i // 3
        col = i % 3
        
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.2 + row * 1.6)
        
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        shape.line.color.rgb = RGBColor(40, 40, 50)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = f"{num}  {name}"
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = indigo_color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(11)
        p2.font.color.rgb = muted_color
        p2.space_before = Pt(4)

    # ═══════════════════════════════════════════════
    # SLIDE 5: INTEGRATION + MARKET
    # ═══════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide5)
    
    add_tag(slide5, "Integration & Market")
    
    title_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Model Agnostic Protection & Market Needs"
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Left column: integration paths
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Integration Schemes"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = emerald_color
    p.space_after = Pt(12)
    
    points = [
        ("🔌 MCP Protocol Proxy", "Redirect your MCP client standard communication channels directly through the Kavach endpoint."),
        ("🐍 SDK Middleware", "Integrate middleware tool execution wrappers directly into LangChain, CrewAI, or LlamaIndex."),
        ("🧹 Pre-Execution Scan", "Pre-process prompts before sending to models to shield against memory leakage and injection.")
    ]
    for header, body in points:
        p_hdr = tf.add_paragraph()
        p_hdr.text = header
        p_hdr.font.size = Pt(16)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = text_color
        p_hdr.space_before = Pt(8)
        
        p_bdy = tf.add_paragraph()
        p_bdy.text = body
        p_bdy.font.size = Pt(12)
        p_bdy.font.color.rgb = muted_color
        p_bdy.space_before = Pt(2)
        
    # Right column: market statistics
    stats = [
        ("$35B+", "AI Security TAM by 2030", indigo_color),
        ("85%", "Enterprises delaying agent adoption", red_color),
        ("3x", "YoY autonomous deployment growth", emerald_color),
        ("0", "Direct runtime agent firewalls today", amber_color)
    ]
    
    for i, (num, label, color) in enumerate(stats):
        row = i // 2
        col = i % 2
        
        x = Inches(6.8 + col * 2.9)
        y = Inches(2.4 + row * 2.2)
        
        shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(1.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        shape.line.color.rgb = RGBColor(40, 40, 50)
        
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        p_num = tf_s.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        p_num.text = num
        p_num.font.name = 'Arial'
        p_num.font.size = Pt(36)
        p_num.font.bold = True
        p_num.font.color.rgb = color
        
        p_lbl = tf_s.add_paragraph()
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.text = label
        p_lbl.font.name = 'Arial'
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = muted_color
        p_lbl.space_before = Pt(6)

    # ═══════════════════════════════════════════════
    # SLIDE 6: ROADMAP
    # ═══════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide6)
    
    add_tag(slide6, "Roadmap")
    
    title_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Roadmap & Future Vision"
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    phases = [
        ("NOW — Q3 2026", "Core Protection", "MCP interception proxy, 9-stage pipeline, SOC dashboard, policy engine, and human approvals.", "✅ BUILT", indigo_color),
        ("Q4 2026", "Advanced Threat Intel", "Vector DB memory poisoning defense. Multi-agent policy consensus. Federated trust networks.", "IN PROGRESS", emerald_color),
        ("Q1 2027", "Self-Learning Platform", "Self-learning anomaly feedback. Real-time sandboxing for compromised agents. SIEM integration.", "PLANNED", amber_color),
        ("Q2 2027", "Enterprise & Scale", "On-premise deployment packages. SOC 2 Type II compliance. Global threat sharing network.", "VISION", purple_color)
    ]
    
    for i, (time, title, desc, badge, color) in enumerate(phases):
        x = Inches(0.8 + i * 2.9)
        shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.7), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(12, 15, 26)
        if badge == "✅ BUILT":
            shape.line.color.rgb = color
            shape.line.width = Pt(1.5)
        else:
            shape.line.color.rgb = RGBColor(40, 40, 50)
            
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = time
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.name = 'Arial'
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = text_color
        p2.space_before = Pt(8)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = 'Arial'
        p3.font.size = Pt(12)
        p3.font.color.rgb = muted_color
        p3.space_before = Pt(8)
        
        p4 = tf.add_paragraph()
        p4.text = badge
        p4.font.name = 'Arial'
        p4.font.size = Pt(11)
        p4.font.bold = True
        p4.font.color.rgb = color
        p4.space_before = Pt(16)

    # ═══════════════════════════════════════════════
    # SLIDE 7: CTA / CONTACT
    # ═══════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide7)
    
    add_tag(slide7, "Let's Build the Future")
    
    title_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Secure the autonomous future."
    p.font.name = 'Arial'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    p2 = tf.add_paragraph()
    p2.text = "The CrowdStrike of Autonomous Agent Systems."
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = indigo_color
    p2.space_before = Pt(6)
    
    # Subtitle explanation
    desc_box = slide7.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(7.5), Inches(1.5))
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]
    p_d.text = "We are building the trust infrastructure that will unlock enterprise adoption of autonomous AI agents. Partner with us to scale agent security."
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(16)
    p_d.font.color.rgb = muted_color
    
    # Info list box
    info_box = slide7.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6.0), Inches(2.2))
    tf_i = info_box.text_frame
    tf_i.word_wrap = True
    
    items = [
        ("📧 Email:", "hello@kavach-ai.dev"),
        ("🌐 Website:", "kavach-ai.dev"),
        ("💻 GitHub:", "github.com/kavach-ai"),
        ("🐦 Twitter:", "@kavach_ai")
    ]
    
    for i, (hdr, val) in enumerate(items):
        p_item = tf_i.paragraphs[0] if i == 0 else tf_i.add_paragraph()
        p_item.text = f"{hdr} {val}"
        p_item.font.size = Pt(14)
        p_item.font.name = 'Arial'
        if i > 0:
            p_item.space_before = Pt(6)
        
        # Color specific segments
        p_item.runs[0].font.bold = True
        p_item.runs[0].font.color.rgb = text_color
        p_item.runs[1].font.color.rgb = purple_color

    # Right side: Metrics breakdown
    metrics_box = slide7.shapes.add_textbox(Inches(8.5), Inches(2.8), Inches(4.0), Inches(4.0))
    tf_m = metrics_box.text_frame
    tf_m.word_wrap = True
    
    metrics = [
        ("50+ Files", "Fully structured, production-ready codebase"),
        ("10K+ Lines", "Comprehensive logic, tests, and API"),
        ("9 Stages", "Zero-trust interception security pipeline"),
        ("<12ms Latency", "Inline intercept decisions & scoring checks")
    ]
    
    for i, (num, desc) in enumerate(metrics):
        p_num = tf_m.paragraphs[0] if i == 0 else tf_m.add_paragraph()
        p_num.text = num
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = emerald_color
        if i > 0:
            p_num.space_before = Pt(12)
            
        p_desc = tf_m.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = muted_color
        p_desc.space_before = Pt(2)
        
    prs.save("pitch-deck.pptx")
    print("Successfully generated pitch-deck.pptx in the workspace root directory!")

if __name__ == "__main__":
    create_pitch_deck()
